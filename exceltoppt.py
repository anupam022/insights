import openai
import pandas as pd
from pptx import Presentation
from pptx.util import Inches

# Set up OpenAI API key
openai.api_key = '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'

# Load Excel data
df = pd.read_excel(f'Book1.xlsx')

# Function to process user query and generate PowerPoint slides
def process_user_query(user_query):
    # Use ChatGPT API to process user query
    response = openai.Completion.create(
        engine="gpt-3.5-turbo",
        prompt=user_query,
        max_tokens=50
    )
    # Extract relevant data from Excel based on user query
    relevant_data = extract_data_from_excel(user_query)
    # Generate PowerPoint slides based on the relevant data
    generate_powerpoint_slides(relevant_data)

# Function to extract relevant data from Excel
def extract_data_from_excel(user_query):
    # Write logic to extract relevant data from Excel based on user query
    return relevant_data

# Function to generate PowerPoint slides
def generate_powerpoint_slides(data):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "Generated PowerPoint Slide"
    content = slide.placeholders[1]
    content.text = "This slide was generated based on the user's query."
    prs.save('output.pptx')

# Example user query
user_query = "What are the sales figures for Q3?"
# Process user query and generate PowerPoint slides
process_user_query(user_query)
